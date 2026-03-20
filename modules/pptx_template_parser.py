"""Gamma生成PPTXのテンプレート解析"""

from pptx import Presentation
from pptx.util import Pt
from modules.models import TemplateTextSlot, TemplateSlotGroup, TemplateMap


def parse_template(pptx_file) -> TemplateMap:
    """PPTXファイルを解析してテンプレート構造を抽出する"""
    prs = Presentation(pptx_file)

    all_slots = []
    slide_titles = []
    slot_groups_per_slide = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_slots = []

        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue

            slot = _extract_slot(slide_idx, shape_idx, shape, text)
            if slot:
                slide_slots.append(slot)
                all_slots.append(slot)

        # スライドタイトルを特定（最大フォントまたは最初のtitle roleのスロット）
        title_slot = None
        for s in slide_slots:
            if s.role == "title":
                title_slot = s
                break
        if not title_slot and slide_slots:
            title_slot = slide_slots[0]
        if title_slot:
            slide_titles.append(title_slot)

        # heading+bodyのグループ化
        groups = _group_slots(slide_idx, slide_slots)
        slot_groups_per_slide.append(groups)

    return TemplateMap(
        source_path="",
        slide_count=len(prs.slides),
        slide_titles=slide_titles,
        slot_groups=slot_groups_per_slide,
        all_slots=all_slots,
    )


def _extract_slot(slide_idx: int, shape_idx: int, shape, text: str) -> TemplateTextSlot | None:
    """シェイプからテキストスロットを抽出"""
    if len(text) < 1:
        return None

    # 最初の段落の最初のランからフォント情報を取得
    font_name = None
    font_size = None
    font_bold = None
    font_color = None

    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.name:
                font_name = run.font.name
            if run.font.size:
                font_size = run.font.size
            if run.font.bold is not None:
                font_bold = run.font.bold
            if run.font.color and run.font.color.rgb:
                font_color = str(run.font.color.rgb)
            break
        if font_name or font_size:
            break

    role = _classify_role(text, font_size, font_bold)

    return TemplateTextSlot(
        slide_index=slide_idx,
        shape_index=shape_idx,
        shape_name=shape.name,
        role=role,
        original_text=text,
        font_name=font_name,
        font_size=font_size,
        font_bold=font_bold,
        font_color_rgb=font_color,
    )


def _classify_role(text: str, font_size: int | None, font_bold: bool | None) -> str:
    """テキストの役割を分類"""
    if font_size:
        pt_size = font_size / 12700  # EMU to pt
        if pt_size >= 28:
            return "title"
        if pt_size >= 16 and font_bold:
            return "heading"

    # フォント情報がない場合はテキスト長で推測
    if len(text) < 30 and font_bold:
        return "heading"
    if len(text) < 20:
        return "heading"

    return "body"


def _group_slots(slide_idx: int, slots: list[TemplateTextSlot]) -> list[TemplateSlotGroup]:
    """heading+bodyの連続ペアをグループ化"""
    groups = []
    i = 0
    while i < len(slots):
        slot = slots[i]
        if slot.role in ("title",):
            i += 1
            continue

        if slot.role == "heading":
            group = TemplateSlotGroup(slide_index=slide_idx, heading_slot=slot)
            # 次がbodyならペアにする
            if i + 1 < len(slots) and slots[i + 1].role == "body":
                group.body_slot = slots[i + 1]
                i += 2
            else:
                i += 1
            groups.append(group)
        elif slot.role == "body":
            # 単独body
            group = TemplateSlotGroup(slide_index=slide_idx, body_slot=slot)
            groups.append(group)
            i += 1
        else:
            i += 1

    return groups
