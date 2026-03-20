"""PowerPointプレゼンテーションの生成"""

import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

from modules.models import PresentationSpec, SlideContent, ParsedTable, ChartSpec
import config

# チャートカラーパレット
CHART_COLORS = [
    RGBColor(0x2E, 0x4A, 0x7A),  # 濃い青
    RGBColor(0xE8, 0x6B, 0x2C),  # オレンジ
    RGBColor(0x4C, 0xA8, 0x5C),  # 緑
    RGBColor(0xC0, 0x39, 0x2B),  # 赤
    RGBColor(0x8E, 0x44, 0xAD),  # 紫
    RGBColor(0xF3, 0x9C, 0x12),  # 黄
    RGBColor(0x16, 0xA0, 0x85),  # ティール
    RGBColor(0x7F, 0x8C, 0x8D),  # グレー
]


def generate_presentation(spec: PresentationSpec, template_path: str | None = None) -> Presentation:
    """PresentationSpecからPowerPointプレゼンテーションを生成する"""
    if template_path:
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    _add_title_slide(prs, spec)

    for slide_content in spec.slides:
        if slide_content.layout == "chart" and slide_content.chart and slide_content.table:
            _add_chart_slide(prs, slide_content)
        elif slide_content.layout == "table_chart" and slide_content.table:
            _add_table_chart_slide(prs, slide_content)
        elif slide_content.layout == "table" and slide_content.table:
            _add_table_slides(prs, slide_content)
        elif slide_content.layout == "bullets":
            _add_bullet_slide(prs, slide_content)
        else:
            _add_text_slide(prs, slide_content)

    return prs


def _set_font(run, size=Pt(14), bold=False, color=None):
    """フォントを設定する共通ヘルパー"""
    run.font.name = config.DEFAULT_FONT
    run.font.size = size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _add_title_slide(prs: Presentation, spec: PresentationSpec):
    """タイトルスライドを追加"""
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)

    # タイトル
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = spec.title
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                _set_font(run, size=Pt(36), bold=True)

    # サブタイトル（対象者・目的）
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = f"対象: {spec.audience}\n目的: {spec.purpose}"
        for paragraph in subtitle.text_frame.paragraphs:
            for run in paragraph.runs:
                _set_font(run, size=Pt(18))


def _add_table_slides(prs: Presentation, slide_content: SlideContent):
    """テーブルスライドを追加（大きいテーブルは分割）"""
    table_data = slide_content.table
    if not table_data:
        return

    max_rows = config.MAX_TABLE_ROWS_PER_SLIDE
    total_rows = len(table_data.rows)

    if total_rows <= max_rows:
        _add_single_table_slide(
            prs, slide_content.title, table_data.headers, table_data.rows,
            slide_content.notes, slide_content.body_text,
        )
    else:
        # テーブルを分割
        for i in range(0, total_rows, max_rows):
            chunk = table_data.rows[i:i + max_rows]
            part_num = i // max_rows + 1
            total_parts = (total_rows + max_rows - 1) // max_rows
            title = f"{slide_content.title} ({part_num}/{total_parts})"
            body = slide_content.body_text if i == 0 else None
            _add_single_table_slide(prs, title, table_data.headers, chunk, slide_content.notes, body)


def _add_single_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list], notes: str = "", body_text: str | None = None):
    """1つのテーブルスライドを追加"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # タイトルをテキストボックスとして追加
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(12)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(28), bold=True)

    # 要約テキストがある場合はテーブルの上に表示
    table_top_offset = Inches(1.3)
    if body_text:
        summary_top = Inches(1.2)
        summary_height = Inches(1.0)
        summary_box = slide.shapes.add_textbox(Inches(0.5), summary_top, Inches(12), summary_height)
        stf = summary_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = body_text
        _set_font(sp.runs[0] if sp.runs else sp.add_run(), size=Pt(12), color=RGBColor(0x33, 0x33, 0x33))
        table_top_offset = Inches(2.4)

    # テーブルを追加
    num_rows = len(rows) + 1  # ヘッダー行含む
    num_cols = len(headers)

    if num_cols == 0:
        return

    table_left = Inches(0.5)
    table_top = table_top_offset
    table_width = Inches(12)
    table_height = Inches(4.5) if body_text else Inches(5.5)

    col_width = table_width // num_cols

    table_shape = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, table_width, table_height)
    table = table_shape.table

    # 列幅を均等に設定
    for i in range(num_cols):
        table.columns[i].width = Emu(col_width)

    # ヘッダー行
    HEADER_BG = RGBColor(0x2E, 0x4A, 0x7A)  # 濃い青
    HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)   # 白
    ALT_ROW_BG = RGBColor(0xE8, 0xEE, 0xF4)  # 薄い青

    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(header)
        # ヘッダーのスタイル
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                _set_font(run, size=Pt(12), bold=True, color=HEADER_FG)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # データ行
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val) if val is not None else ""
            # 交互の背景色
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ALT_ROW_BG
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    _set_font(run, size=Pt(11))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ノート
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _add_bullet_slide(prs: Presentation, slide_content: SlideContent):
    """箇条書きスライドを追加"""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    if slide.shapes.title:
        slide.shapes.title.text = slide_content.title
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            for run in paragraph.runs:
                _set_font(run, size=Pt(28), bold=True)

    if slide_content.body_text and len(slide.placeholders) > 1:
        body = slide.placeholders[1]
        body.text = slide_content.body_text
        for paragraph in body.text_frame.paragraphs:
            for run in paragraph.runs:
                _set_font(run, size=Pt(16))

    if slide_content.notes:
        slide.notes_slide.notes_text_frame.text = slide_content.notes


def _add_text_slide(prs: Presentation, slide_content: SlideContent):
    """テキストスライドを追加"""
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = slide_content.title
    _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(28), bold=True)

    # 本文
    if slide_content.body_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = slide_content.body_text
        _set_font(p2.runs[0] if p2.runs else p2.add_run(), size=Pt(16))


def _to_number(val) -> float:
    """値を数値に変換（変換不可の場合は0）"""
    if isinstance(val, (int, float)) and not isinstance(val, bool):
        return float(val)
    if isinstance(val, str):
        try:
            return float(val.replace(",", "").replace("%", ""))
        except (ValueError, AttributeError):
            return 0.0
    return 0.0


def _get_chart_type_enum(chart_type: str) -> XL_CHART_TYPE:
    """文字列からpptxのチャートタイプに変換"""
    mapping = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
    }
    return mapping.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)


def _build_chart_data(table: ParsedTable, chart_spec: ChartSpec) -> CategoryChartData:
    """テーブルデータからチャートデータを構築"""
    chart_data = CategoryChartData()

    # カテゴリ（ラベル）を設定
    cat_col = chart_spec.category_column
    categories = []
    for row in table.rows:
        if cat_col < len(row) and row[cat_col] is not None:
            categories.append(str(row[cat_col]))
        else:
            categories.append("")
    chart_data.categories = categories

    # 値系列を追加
    for val_col in chart_spec.value_columns:
        series_name = table.headers[val_col] if val_col < len(table.headers) else f"系列{val_col}"
        values = []
        for row in table.rows:
            if val_col < len(row):
                values.append(_to_number(row[val_col]))
            else:
                values.append(0.0)
        chart_data.add_series(series_name, values)

    return chart_data


def _style_chart(chart, chart_spec: ChartSpec):
    """チャートのスタイルを設定"""
    chart.has_legend = len(chart_spec.value_columns) > 1 or chart_spec.chart_type == "pie"
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.legend.font.name = config.DEFAULT_FONT

    # 円グラフの場合はデータラベルを表示
    if chart_spec.chart_type == "pie":
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(10)
        data_labels.font.name = config.DEFAULT_FONT
        data_labels.number_format = "0.0%"
        data_labels.show_percentage = True
        data_labels.show_category_name = True
        data_labels.show_value = False

    # 系列の色を設定
    for i, series in enumerate(chart.series):
        color = CHART_COLORS[i % len(CHART_COLORS)]
        if chart_spec.chart_type == "pie":
            # 円グラフはポイントごとに色を設定
            for j in range(len(series.values)):
                point = series.points[j]
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = CHART_COLORS[j % len(CHART_COLORS)]
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
            if chart_spec.chart_type == "line":
                series.format.line.color.rgb = color
                series.format.line.width = Pt(2.5)
                series.smooth = False

    # 軸のフォント設定（円グラフ以外）
    if chart_spec.chart_type != "pie":
        if chart.category_axis:
            chart.category_axis.tick_labels.font.size = Pt(9)
            chart.category_axis.tick_labels.font.name = config.DEFAULT_FONT
        if chart.value_axis:
            chart.value_axis.tick_labels.font.size = Pt(9)
            chart.value_axis.tick_labels.font.name = config.DEFAULT_FONT


def _add_chart_slide(prs: Presentation, slide_content: SlideContent):
    """チャートスライドを追加"""
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = slide_content.title
    _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(28), bold=True)

    # 要約テキスト
    chart_top = Inches(1.3)
    if slide_content.body_text:
        summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.8))
        stf = summary_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = slide_content.body_text
        _set_font(sp.runs[0] if sp.runs else sp.add_run(), size=Pt(12), color=RGBColor(0x33, 0x33, 0x33))
        chart_top = Inches(2.2)

    # チャートを追加
    chart_data = _build_chart_data(slide_content.table, slide_content.chart)
    chart_type = _get_chart_type_enum(slide_content.chart.chart_type)
    chart_height = Inches(7.0 - chart_top.inches)

    chart_frame = slide.shapes.add_chart(
        chart_type,
        Inches(0.8), chart_top, Inches(11.5), chart_height,
        chart_data,
    )
    _style_chart(chart_frame.chart, slide_content.chart)

    # ノート
    if slide_content.notes:
        slide.notes_slide.notes_text_frame.text = slide_content.notes


def _add_table_chart_slide(prs: Presentation, slide_content: SlideContent):
    """テーブル+チャートの並列スライドを追加"""
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = slide_content.title
    _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(28), bold=True)

    # 要約テキスト
    content_top = Inches(1.3)
    if slide_content.body_text:
        summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.7))
        stf = summary_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = slide_content.body_text
        _set_font(sp.runs[0] if sp.runs else sp.add_run(), size=Pt(11), color=RGBColor(0x33, 0x33, 0x33))
        content_top = Inches(2.1)

    table_data = slide_content.table
    content_height = Inches(7.2 - content_top.inches)

    # 左側: テーブル（コンパクト）
    num_rows = min(len(table_data.rows), 12) + 1
    num_cols = len(table_data.headers)
    if num_cols > 0:
        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(0.3), content_top, Inches(6.0), content_height,
        )
        table = table_shape.table

        HEADER_BG = RGBColor(0x2E, 0x4A, 0x7A)
        HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
        ALT_ROW_BG = RGBColor(0xE8, 0xEE, 0xF4)

        col_width = Emu(Inches(6.0) // num_cols)
        for i in range(num_cols):
            table.columns[i].width = col_width

        for j, header in enumerate(table_data.headers):
            cell = table.cell(0, j)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = HEADER_BG
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    _set_font(run, size=Pt(9), bold=True, color=HEADER_FG)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        display_rows = table_data.rows[:12]
        for i, row in enumerate(display_rows):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(val) if val is not None else ""
                if i % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = ALT_ROW_BG
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        _set_font(run, size=Pt(9))
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 右側: チャート
    if slide_content.chart:
        chart_data = _build_chart_data(table_data, slide_content.chart)
        chart_type = _get_chart_type_enum(slide_content.chart.chart_type)

        chart_frame = slide.shapes.add_chart(
            chart_type,
            Inches(6.5), content_top, Inches(6.3), content_height,
            chart_data,
        )
        _style_chart(chart_frame.chart, slide_content.chart)

    # ノート
    if slide_content.notes:
        slide.notes_slide.notes_text_frame.text = slide_content.notes


def to_bytes(prs: Presentation) -> bytes:
    """プレゼンテーションをバイト列に変換"""
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
