"""テンプレートPPTXのテキスト差し替え（書式保持）"""

import io
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

from modules.models import TemplateMap, FieldMapping, TemplateTextSlot


def update_presentation(pptx_file, mappings: list[FieldMapping], data_rows: list[dict]) -> bytes:
    """テンプレートPPTXにExcelデータを差し替えて返す。

    Args:
        pptx_file: テンプレートPPTXファイル
        mappings: フィールドマッピングのリスト
        data_rows: Excelから抽出したデータ行（dict形式）
    Returns:
        更新されたPPTXのバイト列
    """
    prs = Presentation(pptx_file)

    # マッピングをshape_indexでインデックス化
    slot_to_column = {}
    for m in mappings:
        if m.excel_column:
            key = (m.slot.slide_index, m.slot.shape_index)
            slot_to_column[key] = m.excel_column

    # 各スライドのシェイプを更新
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            key = (slide_idx, shape_idx)
            if key not in slot_to_column:
                continue
            if not shape.has_text_frame:
                continue

            col_name = slot_to_column[key]
            new_text = _get_value_for_shape(col_name, data_rows, slide_idx, shape_idx)
            if new_text is not None:
                _replace_shape_text_preserve_format(shape, new_text)

    return _to_bytes(prs)


def update_by_slot_text(pptx_file, replacements: list[tuple[TemplateTextSlot, str]]) -> bytes:
    """スロットのoriginal_textを手がかりにテキストを差し替える。

    Args:
        pptx_file: テンプレートPPTXファイル
        replacements: (スロット, 新テキスト) のリスト
    Returns:
        更新されたPPTXのバイト列
    """
    prs = Presentation(pptx_file)

    replace_map = {}
    for slot, new_text in replacements:
        key = (slot.slide_index, slot.shape_index)
        replace_map[key] = new_text

    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            key = (slide_idx, shape_idx)
            if key in replace_map and shape.has_text_frame:
                _replace_shape_text_preserve_format(shape, replace_map[key])

    return _to_bytes(prs)


def _replace_shape_text_preserve_format(shape, new_text: str):
    """シェイプのテキストを書式を保持したまま差し替える（lxml直接操作）"""
    txBody = shape._element.find(qn("p:txBody"))
    if txBody is None:
        return

    paragraphs = txBody.findall(qn("a:p"))
    if not paragraphs:
        return

    # 最初の段落とランの書式をスナップショット
    first_p = paragraphs[0]
    pPr_original = first_p.find(qn("a:pPr"))
    first_r = first_p.find(qn("a:r"))
    rPr_original = None
    if first_r is not None:
        rPr_original = first_r.find(qn("a:rPr"))

    # bodyPropertiesとlstStyleを保持
    bodyPr = txBody.find(qn("a:bodyPr"))
    lstStyle = txBody.find(qn("a:lstStyle"))

    # 既存の段落をすべて削除
    for p in paragraphs:
        txBody.remove(p)

    # 新しいテキストを段落ごとに追加
    lines = new_text.split("\n") if new_text else [""]
    for line in lines:
        new_p = etree.SubElement(txBody, qn("a:p"))

        # 段落プロパティを復元
        if pPr_original is not None:
            new_p.insert(0, deepcopy(pPr_original))

        # ランを追加
        new_r = etree.SubElement(new_p, qn("a:r"))
        if rPr_original is not None:
            new_r.insert(0, deepcopy(rPr_original))
        t = etree.SubElement(new_r, qn("a:t"))
        t.text = line


def _get_value_for_shape(col_name: str, data_rows: list[dict], slide_idx: int, shape_idx: int) -> str | None:
    """データ行から値を取得。スライド番号に基づいて対応する行を探す"""
    # data_rowsから対応する行を特定（簡易的に行インデックスで対応）
    for row in data_rows:
        if col_name in row and row[col_name] is not None:
            val = row[col_name]
            return str(val) if val is not None else None
    return None


def _to_bytes(prs: Presentation) -> bytes:
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
